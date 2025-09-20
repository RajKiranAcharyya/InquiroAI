import fitz  # PyMuPDF
from io import BytesIO
from docx import Document
from pptx import Presentation

def extract_text_from_docx(file):
    doc = Document(BytesIO(file.read()))
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return "\n".join(full_text)

def extract_text_from_pptx(file):
    presentation = Presentation(BytesIO(file.read()))
    full_text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return "\n".join(full_text)

def extract_text_from_files(files):
    """Extract text from a list of file-like objects (pdf, docx, pptx) and return as a list of documents."""
    documents = []
    for file in files:
        filename = file.name.lower()
        file.seek(0)  # Reset file pointer before reading
        if filename.endswith(".pdf"):
            doc = fitz.open(stream=file.read(), filetype="pdf")
            full_text = ""
            for page in doc:
                full_text += page.get_text()
            documents.append(full_text)
        elif filename.endswith(".docx"):
            file.seek(0)
            documents.append(extract_text_from_docx(file))
        elif filename.endswith(".pptx"):
            file.seek(0)
            documents.append(extract_text_from_pptx(file))
        else:
            documents.append("")  # Unsupported file type
    return documents




# import fitz 

# def extract_text_from_pdfs(pdf_files):
#     """Extract text from a list of PDF file-like objects and return as a list of documents."""
#     documents = []
#     for pdf_file in pdf_files:
#         doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
#         full_text = ""
#         for page in doc:
#             full_text += page.get_text()
#         documents.append(full_text)
#     return documents

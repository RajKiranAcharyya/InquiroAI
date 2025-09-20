import streamlit as st
from io import BytesIO
import fitz  # PyMuPDF for PDFs
from docx import Document
from pptx import Presentation
from qa_engine import AcademicQAEngine

# CSS styling
st.set_page_config(page_title="Inquiro", layout="wide")
st.markdown("""
    <style>
    .main {
        padding: 2rem;
    }
    .stTitle {
        color: #2E4057;
        font-size: 3rem !important;
        text-align: center;
        padding-bottom: 2rem;
    }
    .stTextInput > div > div > input {
        background-color: #f0f2f6;
        padding: 1rem;
        font-size: 1.1rem;
        border-radius: 10px;
    }
    .chat-message {
        padding: 1.5rem;
        border-radius: 10px;
        margin-bottom: 1rem;
        background-color: #f0f2f6;
    }
    .user-message {
        background-color: #e1e9f7;
    }
    .sidebar .element-container {
        background-color: #ffffff;
        padding: 1rem;
        border-radius: 10px;
        margin-bottom: 1rem;
    }
    </style>
    """, unsafe_allow_html=True)

st.title("🎓 Inquiro Smart AI Knowledge Hub")
st.markdown("<p style='text-align: center; color: #666666; font-size: 1.2rem;'>Your AI-powered research companion</p>", unsafe_allow_html=True)

# Sidebar for multi-format file upload
with st.sidebar:
    st.markdown("### 📚 Document Upload")
    uploaded_files = st.file_uploader(
        "Select PDF, PPTX or DOCX files",
        accept_multiple_files=True,
        type=["pdf", "pptx", "docx"],
        help="Upload one or more PDF, PPTX, or DOCX documents to begin"
    )
    if uploaded_files:
        st.success(f"📎 {len(uploaded_files)} file(s) uploaded successfully")
    else:
        st.info("👆 Please upload files to start")

# Extraction functions
def extract_text_from_docx(file):
    doc = Document(BytesIO(file.read()))
    full_text = [para.text for para in doc.paragraphs]
    return "\n".join(full_text)

def extract_text_from_pptx(file):
    presentation = Presentation(BytesIO(file.read()))
    full_text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return "\n".join(full_text)

def extract_text_from_files(files):
    docs = []
    for file in files:
        filename = file.name.lower()
        file.seek(0)
        if filename.endswith(".pdf"):
            pdf_doc = fitz.open(stream=file.read(), filetype="pdf")
            text = ""
            for page in pdf_doc:
                text += page.get_text()
            docs.append(text)
        elif filename.endswith(".docx"):
            file.seek(0)
            docs.append(extract_text_from_docx(file))
        elif filename.endswith(".pptx"):
            file.seek(0)
            docs.append(extract_text_from_pptx(file))
        else:
            docs.append("")
    return docs

# Session state initialization
if "uploaded_files" not in st.session_state:
    st.session_state.uploaded_files = None
if "documents" not in st.session_state:
    st.session_state.documents = None
if "qa_engine" not in st.session_state:
    st.session_state.qa_engine = None
if "questions" not in st.session_state:
    st.session_state.questions = []
if "answers" not in st.session_state:
    st.session_state.answers = []
if "clear_input" not in st.session_state:
    st.session_state.clear_input = False

# Store uploaded files persistently
if uploaded_files:
    st.session_state.uploaded_files = uploaded_files

# Initialize docs and QA engine once
if (
    st.session_state.uploaded_files
    and (st.session_state.documents is None or st.session_state.qa_engine is None)
):
    with st.spinner("📑 Processing your documents..."):
        st.session_state.documents = extract_text_from_files(st.session_state.uploaded_files)
        st.session_state.qa_engine = AcademicQAEngine(st.session_state.documents)
    st.success("✨ Documents processed successfully! Ask away!")
    st.session_state.questions = []
    st.session_state.answers = []
    st.session_state.clear_input = False

# Clear input before rendering widget if needed
if st.session_state.clear_input:
    st.session_state.user_question_input = ""
    st.session_state.clear_input = False

# Main chat interface
if st.session_state.qa_engine:
    st.markdown("### 💭 Ask Your Question")
    user_question = st.text_input(
        label="",
        placeholder="Type your question here and press Enter...",
        key="user_question_input",
    )

    if user_question:
        with st.spinner("🤔 Thinking..."):
            answer = st.session_state.qa_engine.answer_question(user_question)
        st.session_state.questions.append(user_question)
        st.session_state.answers.append(answer)
        st.session_state.clear_input = True

    if st.session_state.questions:
        st.markdown("### 📝 Conversation History")
        for q, a in zip(reversed(st.session_state.questions), reversed(st.session_state.answers)):
            st.markdown(f"""
                <div class="chat-message user-message">
                    <strong>You:</strong> {q}
                </div>
                <div class="chat-message">
                    <strong>Assistant:</strong> {a}
                </div>
                """, unsafe_allow_html=True)
else:
    st.markdown("""
        <div style='text-align: center; padding: 3rem;'>
            <h3>👋 Welcome!</h3>
            <p>Upload your documents in the sidebar to start exploring them with AI assistance.</p>
        </div>
        """, unsafe_allow_html=True)














# import streamlit as st
# from pdf_utils import extract_text_from_pdfs
# from qa_engine import AcademicQAEngine

# st.set_page_config(page_title="Inquiro", layout="wide")
# st.markdown("""
#     <style>
#     .main {
#         padding: 2rem;
#     }
#     .stTitle {
#         color: #2E4057;
#         font-size: 3rem !important;
#         text-align: center;
#         padding-bottom: 2rem;
#     }
#     .stTextInput > div > div > input {
#         background-color: #f0f2f6;
#         padding: 1rem;
#         font-size: 1.1rem;
#         border-radius: 10px;
#     }
#     .chat-message {
#         padding: 1.5rem;
#         border-radius: 10px;
#         margin-bottom: 1rem;
#         background-color: #f0f2f6;
#     }
#     .user-message {
#         background-color: #e1e9f7;
#     }
#     .sidebar .element-container {
#         background-color: #ffffff;
#         padding: 1rem;
#         border-radius: 10px;
#         margin-bottom: 1rem;
#     }
#     </style>
#     """, unsafe_allow_html=True)

# st.title("🎓 Inquiro Smart AI Knowlwdge Hub")
# st.markdown("<p style='text-align: center; color: #666666; font-size: 1.2rem;'>Your AI-powered research companion</p>", unsafe_allow_html=True)

# with st.sidebar:
#     st.markdown("### 📚 Document Upload")
#     uploaded_files = st.file_uploader(
#         "Select PDF files", accept_multiple_files=True, type=["pdf"],
#         help="Upload one or more PDF documents to begin"
#     )
#     if uploaded_files:
#         st.success(f"📎 {len(uploaded_files)} file(s) uploaded successfully")
#     else:
#         st.info("👆 Please upload PDF files to start")

# # Session state initialization
# if "uploaded_files" not in st.session_state:
#     st.session_state.uploaded_files = None
# if "documents" not in st.session_state:
#     st.session_state.documents = None
# if "qa_engine" not in st.session_state:
#     st.session_state.qa_engine = None
# if "questions" not in st.session_state:
#     st.session_state.questions = []
# if "answers" not in st.session_state:
#     st.session_state.answers = []
# if "clear_input" not in st.session_state:
#     st.session_state.clear_input = False

# if uploaded_files:
#     st.session_state.uploaded_files = uploaded_files

# if (
#     st.session_state.uploaded_files
#     and (st.session_state.documents is None or st.session_state.qa_engine is None)
# ):
#     with st.spinner("📑 Processing your documents..."):
#         st.session_state.documents = extract_text_from_pdfs(st.session_state.uploaded_files)
#         st.session_state.qa_engine = AcademicQAEngine(st.session_state.documents)
#     st.success("✨ Documents processed successfully! Ask away!")
#     st.session_state.questions = []
#     st.session_state.answers = []
#     st.session_state.clear_input = False

# if st.session_state.clear_input:
#     st.session_state.user_question_input = ""
#     st.session_state.clear_input = False

# if st.session_state.qa_engine:
#     st.markdown("### 💭 Ask Your Question")
#     user_question = st.text_input(
#         label="",
#         placeholder="Type your question here and press Enter...",
#         key="user_question_input",
#     )

#     if user_question:
#         with st.spinner("🤔 Thinking..."):
#             answer = st.session_state.qa_engine.answer_question(user_question)
#         st.session_state.questions.append(user_question)
#         st.session_state.answers.append(answer)
#         st.session_state.clear_input = True

#     if st.session_state.questions:
#         st.markdown("### 📝 Conversation History")
#         for q, a in zip(reversed(st.session_state.questions), reversed(st.session_state.answers)):
#             st.markdown(f"""
#                 <div class="chat-message user-message">
#                     <strong>You:</strong> {q}
#                 </div>
#                 <div class="chat-message">
#                     <strong>Assistant:</strong> {a}
#                 </div>
#                 """, unsafe_allow_html=True)
# else:
#     st.markdown("""
#         <div style='text-align: center; padding: 3rem;'>
#             <h3>👋 Welcome!</h3>
#             <p>Upload your PDF documents in the sidebar to start exploring them with AI assistance.</p>
#         </div>
#         """, unsafe_allow_html=True)
